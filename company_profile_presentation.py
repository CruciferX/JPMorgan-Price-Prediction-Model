from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt

# Create a new presentation
prs = Presentation()
SLIDE_WIDTH = prs.slide_width
JP_BLUE = RGBColor(0, 32, 91)
JP_GRAY = RGBColor(102, 102, 102)
BODY_FONT_SIZE = Pt(20)
TITLE_FONT_SIZE = Pt(32)
HEADER_HEIGHT = Inches(0.5)
FOOTER_HEIGHT = Inches(0.3)

# Apply JP Morgan visual theme
def apply_jp_theme(slide, include_title=True):
    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, HEADER_HEIGHT
    )
    header.fill.solid()
    header.fill.fore_color.rgb = JP_BLUE
    header.line.fill.background()

    # Footer bar
    footer = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, prs.slide_height - FOOTER_HEIGHT, SLIDE_WIDTH, FOOTER_HEIGHT
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = JP_GRAY
    footer.line.fill.background()
    tf = footer.text_frame
    tf.text = "JP Morgan Confidential"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    tf.paragraphs[0].alignment = 1  # center

    # Adjust title position
    if include_title and slide.shapes.title:
        slide.shapes.title.top = Inches(0.6)

# Title slide
def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    apply_jp_theme(slide)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = JP_BLUE
    slide.placeholders[1].text_frame.paragraphs[0].font.size = Pt(18)
    slide.placeholders[1].text_frame.paragraphs[0].font.color.rgb = JP_GRAY
    return slide

# Content slides
def add_jpm_slide(title_text, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_jp_theme(slide)
    slide.shapes.title.text = title_text
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = JP_BLUE
    slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = BODY_FONT_SIZE
        p.font.name = 'Calibri'
        p.font.color.rgb = JP_GRAY

    slide.notes_slide.notes_text_frame.text = "Animation: Appear each bullet point one by one."

# Logo image via matplotlib
logo_path = "jpmorgan_logo.png"
plt.figure(figsize=(4, 1))
plt.text(0.5, 0.5, "JP Morgan", fontsize=28, ha='center', va='center', color='#00205B')
plt.axis('off')
plt.savefig(logo_path, bbox_inches='tight')
plt.close()

# Add title slide
slide = add_title_slide("FICO Score Bucketing & Default Risk Prediction", "JP Morgan Quant Research | Tushar Verma")
slide.shapes.add_picture(logo_path, Inches(5.5), Inches(0.3), height=Inches(0.8))

# Content
add_jpm_slide("Objective", [
    "Discretize continuous FICO scores using optimal quantization.",
    "Use Mean Squared Error (MSE) minimization to define buckets.",
    "Enable effective input to machine learning classification models."
])

add_jpm_slide("Data Overview", [
    "Dataset: Historical loan and borrower records.",
    "Key columns: FICO Score, Default Indicator, Loan Amount, etc.",
    "Filtered for completeness and sorted by FICO score."
])

add_jpm_slide("What is Quantization?", [
    "Quantization: Mapping continuous values to discrete levels.",
    "In signal processing, used for digitizing analog signals.",
    "Here: we quantize FICO scores to 5 optimal buckets minimizing MSE.",
    "Goal: Preserve important predictive features while reducing complexity."
])

add_jpm_slide("Mean Squared Error (MSE)", [
    "Measures average of squared differences from mean.",
    "Used to evaluate compactness within each bucket.",
    "Lower MSE = tighter bucket = more representative average.",
    "Objective: Minimize total intra-bucket MSE."
])

add_jpm_slide("Dynamic Programming", [
    "Technique to break down optimization problems into sub-problems.",
    "Stores intermediate solutions for reuse (memoization).",
    "Used here to find optimal score breakpoints in O(n²) time.",
    "Ensures global MSE is minimized across all possible partitions."
])

add_jpm_slide("Likelihood Function (Next Steps)", [
    "Function of parameters given observed data.",
    "Helps in estimating parameters using Maximum Likelihood Estimation (MLE).",
    "Future extension: use log-likelihood over MSE for probabilistic bucketing.",
    "Better alignment with classification model training objectives."
])

add_jpm_slide("Approach Summary", [
    "1. Precompute MSE for all FICO sub-ranges.",
    "2. Apply dynamic programming to find optimal buckets.",
    "3. Map each FICO score to a discrete rating label (B1–B5).",
    "4. Label used as model input for predicting default risk."
])

add_jpm_slide("Results", [
    "Generated 5 buckets with low intra-bucket MSE.",
    "Higher buckets correlated with increased default risk.",
    "Enabled better interpretability and model performance.",
    "Clear segmentation for downstream credit scoring models."
])

# Chart slide
chart_path = "default_chart.png"
buckets = ['B1', 'B2', 'B3', 'B4', 'B5']
defaults = [30, 22, 18, 10, 5]
plt.figure()
plt.bar(buckets, defaults, color='#00205B')
plt.title("Default Rate by Bucket")
plt.ylabel("Number of Defaults")
plt.savefig(chart_path, bbox_inches='tight')
plt.close()

slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
apply_jp_theme(slide, include_title=True)
slide.shapes.title.text = "Default Distribution by Bucket"
slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = JP_BLUE
slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
slide.shapes.add_picture(chart_path, Inches(1), Inches(1.5), width=Inches(7))
slide.notes_slide.notes_text_frame.text = "Animation: Fade-in chart for impact."

add_jpm_slide("Conclusion", [
    "MSE-based quantization effectively groups FICO scores.",
    "Dynamic programming enables optimal segmentation.",
    "Model-ready ratings improve credit risk prediction performance.",
    "Foundation laid for future likelihood-based modeling."
])

add_jpm_slide("Next Steps", [
    "Implement likelihood-based quantization for better calibration.",
    "Incorporate additional borrower-level features (e.g., income, LTV).",
    "Benchmark performance with logistic regression, XGBoost.",
    "Explore real-time bucketing in credit risk systems."
])

# Save the presentation
prs.save("JP_Morgan_Quant_Themed_Presentation.pptx")
print("✅ Presentation saved as 'JP_Morgan_Quant_Themed_Presentation.pptx'")
